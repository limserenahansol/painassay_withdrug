"""fix_hansol_deck.py  -  correct Hansol's deck into a NEW file.

Never writes to the input. Keeps every layout choice - positions, sizes,
fonts - and changes only what the audit found wrong:

  1  "Day 1 Vehicle" was CORRECT and must be left alone. An earlier version
     of this script changed it to "no drug" because the randomisation sheet
     labels rows 1-6 "Baseline (no injection)" with Treatment = None.
     Hansol ran vehicle on Day 1. The sheet does not describe what was done.
     Consequence: the injection itself IS controlled, and the remaining
     confound is only the fixed order - vehicle always first, SBI always
     second - not the absence of a vehicle arm.
  2  "fell to 8 %"      -> 22 %   (8 % was the response-window value)
  3  "withdrawal to 72 %" -> 63 %
  4  slide 7 holds an OLDER S1 (1800x978 px; current is 1800x1027) whose own
     title still reads "fell the MOST, not the least". Replaced.
  5  that image was placed at aspect 1.67 against its true 1.84, a 9 %
     horizontal squash. Re-fitted.
  6  a sentence on slide 2 ends in a comma. Completed.

Also shortens the wordiest blocks and replaces titles with one plain clause
each, since two slides had near-identical titles.

USAGE
    python fix_hansol_deck.py                 # writes *_checked.pptx
    python fix_hansol_deck.py --dry-run       # report only, write nothing

Hansol Lim - HEAL mini1p / SBI-553
"""
from __future__ import annotations

import argparse
import os
import shutil

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

BASE = r"C:\Users\hsollim\Documents\HEAL_mini1p_SBI553"
SRC = os.path.join(BASE, "lab_meeting", "WT_SBI553_FINAL_hansol.pptx")
S1 = os.path.join(BASE, "lab_meeting", "figs_day2", "S1_sedation_evidence.png")

# exact-substring replacements, applied run by run so fonts survive
TEXT_FIXES = [
    # --- factual ---
    # NOTE: "Day 1 Vehicle" is correct - do not touch it. See item 1 above.
    ("fell to 8 % of Day 1 at day 2 — the largest drop of any behaviour. "
     "Reflexes fell least",
     "fell to 22 % of Day 1 — as far as licking/biting (20 %), the joint "
     "largest drop. Reflexes fell least"),
    ("(withdrawal to 72 %).", "(withdrawal to 63 %, flinch to 82 %)."),
    ("All six mice got SBI-553, ",
     "Vehicle was given on Day 1 and SBI-553 on Day 2, so the injection is "
     "controlled but the fixed order is not."),
    # --- titles, one plain clause each ---
    ("Day 2 result, in one line",
     "The drug reduced everything, not just pain"),
    ("Session design", "How one session runs"),
    ("Six behaviours, two classes - kept separate on purpose",
     "What we score, and why we keep two lists"),
    ("What the scoring/Video looks like",
     "What one scored session looks like"),
    ("Everything fell ", "Every mouse, every behaviour"),
    ("Every behaviour fell — including escape/rearing",
     "Exploration fell as much as pain"),
    ("Lower on the drug day for all six behaviours, and flatter for the "
     "affective ones", "Lower on the drug day, every behaviour"),
    ("Paired mouse by mouse: paw withdrawal and escape/rearing decrease "
     "mainly", "Same mouse, both days"),
    ("The baseline is the decisive comparison",
     "Already slower before any stimulus"),
    # --- trim the wordiest blocks ---
    ("METHOD   Change index per mouse, dot = the index, line = 95 % CI from "
     "the Poisson variance of the totals. EXACT POISSON RATE TEST on that "
     "animal's own counts: conditional on the n₁+n₂ events in total, the "
     "number landing on Day 2 is Binomial(n₁+n₂, t₂/(t₁+t₂)) if the two "
     "rates are equal; the two-sided binomial p is exact. Stars are "
     "uncorrected — 25 of 36 tests reach p<0.05, 22 survive "
     "Benjamini-Hochberg.",
     "METHOD   One dot per mouse, line = 95 % CI. Exact Poisson rate test on "
     "that animal's own counts (57–100 stimuli each). Events counted inside "
     "the four stimulus blocks. *** p<0.001  ** p<0.01  * p<0.05.  Stars "
     "uncorrected: 25 of 36 reach p<0.05, 22 survive Benjamini-Hochberg."),
    ("Medians: licking/biting ×0.20 and escape/rearing ×0.22 fell most, paw "
     "withdrawal ×0.63 and flinch ×0.82 least (attending ×0.51, guarding "
     "×0.45). Escape/rearing is exploration — Day 1 showed it is high at "
     "baseline and DECREASES under stimulation — so a selective analgesic "
     "should have spared it. It did not.",
     "Medians: licking/biting ×0.20 and escape/rearing ×0.22 fell most; "
     "withdrawal ×0.63 and flinch ×0.82 least. Escape/rearing is "
     "exploration, not pain — a selective analgesic should have spared it."),
    ("Paw withdrawal looks close to the line only because the axis is "
     "logarithmic: ×0.63 is a 37 % decrease, in all six mice, paired "
     "p = 0.031. It is significantly reduced AND the least reduced of the "
     "six.",
     "Paw withdrawal sits near the line only because the axis is "
     "logarithmic: ×0.63 is a 37 % fall, in all six mice, p = 0.031."),
    ("Escape/rearing is significant at all four stimuli. Paw withdrawal is "
     "significant at three of four. Guarding and flinch are not significant "
     "per stimulus — with six mice the test cannot go below 0.031",
     "Escape/rearing is significant at all four stimuli, withdrawal at three "
     "of four. With six mice the test cannot go below p = 0.031, so a blank "
     "cell means underpowered, not unchanged."),
    ("A selective affective effect would move every arrow straight DOWN. F2, "
     "F3 and M3 do move mostly down; M1 and M2 move down AND left; F1 moves "
     "left and slightly up. Combined with escape/rearing falling hardest, "
     "the pattern is a general reduction in activity.",
     "A selective affective effect would move every arrow straight DOWN. "
     "Most move down AND left, so both classes fell together — a general "
     "reduction in activity."),
    ("Escape/rearing decreased in all 6 mice (6/6 significant). "
     "Licking/biting, attending and paw withdrawal also fell in 6/6. "
     "Reflexes fell least. One mouse increased guarding 9-fold.",
     "Escape/rearing fell in all 6 mice, the only behaviour significant in "
     "every one. Licking/biting, attending and withdrawal also fell in 6/6. "
     "Reflexes fell least. One mouse (F1) increased guarding 9-fold."),
    ("same six mice, dosed 10 min before the assay. ",
     "same six mice.  SBI-553 given 10 min before the Day 2 assay."),
]


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--src", default=SRC)
    ap.add_argument("--out", default=None)
    ap.add_argument("--dry-run", action="store_true")
    a = ap.parse_args()
    out = a.out or a.src.replace(".pptx", "_checked.pptx")
    if os.path.abspath(out) == os.path.abspath(a.src):
        raise SystemExit("refusing to overwrite the original")

    if not a.dry_run:
        shutil.copy2(a.src, out)
    prs = Presentation(out if not a.dry_run else a.src)

    # ---- text ----
    # Two passes. Per-run first, which preserves any formatting that varies
    # inside a paragraph. Then per-paragraph for the patterns that a run
    # boundary cut in half - PowerPoint splits a run wherever the user
    # retyped or the spell-checker touched it, so a plain substring test on
    # runs misses roughly a third of them. The paragraph pass writes the
    # whole string into the first run and empties the others, which is safe
    # here because these blocks are a single size and colour throughout.
    done, missed = [], list(TEXT_FIXES)

    def note(i, old, new, how):
        done.append((i, old[:52], how))
        if (old, new) in missed:
            missed.remove((old, new))

    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    for old, new in TEXT_FIXES:
                        if old and old in r.text:
                            r.text = r.text.replace(old, new)
                            note(i, old, new, "run")
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                runs = para.runs
                if len(runs) < 2:
                    continue
                joined = "".join(r.text for r in runs)
                hit = False
                for old, new in TEXT_FIXES:
                    if old and old in joined:
                        joined = joined.replace(old, new)
                        note(i, old, new, "paragraph")
                        hit = True
                if hit:
                    runs[0].text = joined
                    for r in runs[1:]:
                        r.text = ""
    print(f"text: {len(done)} replacement(s)")
    for i, o, how in done:
        print(f"  slide {i:2d}  [{how:9s}] {o}")
    if missed:
        print(f"\n  {len(missed)} pattern(s) NOT found - the run was split "
              f"across formatting, so these need doing by hand:")
        for o, _ in missed:
            print(f"    {o[:70]}")

    # ---- the stale, squashed S1 on slide 7 ----
    if os.path.exists(S1):
        iw, ih = Image.open(S1).size
        asp = iw / ih
        for i, sl in enumerate(prs.slides, 1):
            for sh in list(sl.shapes):
                if sh.__class__.__name__ != "Picture":
                    continue
                try:
                    pw, ph = Image.open(__import__("io").BytesIO(
                        sh.image.blob)).size
                except Exception:
                    continue
                # the old S1: same width, shorter, and not the current file
                if pw == 1800 and ph != ih:
                    L, T = sh.left, sh.top
                    W = sh.width
                    H = Inches(W.inches / asp)
                    print(f"\nslide {i}: replacing the stale S1 "
                          f"({pw}x{ph} px) with the current one "
                          f"({iw}x{ih} px), aspect corrected "
                          f"{sh.width.inches / sh.height.inches:.2f} "
                          f"-> {asp:.2f}")
                    if not a.dry_run:
                        sh._element.getparent().remove(sh._element)
                        sl.shapes.add_picture(S1, L, T, W, H)

    if a.dry_run:
        print("\n--dry-run: nothing written")
        return
    prs.save(out)
    print(f"\nwrote {out}")
    print(f"original untouched: {a.src}")


if __name__ == "__main__":
    main()
