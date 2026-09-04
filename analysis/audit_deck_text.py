"""READ-ONLY audit of the final deck, scope-aware. Changes nothing.

The previous version failed five claims because it compared every number
against the block-scope CSV. Two figures legitimately use different scopes:

  S1 (slide 1)                 every event in the WHOLE SESSION
  forest / D2 / D8 (2,3,4)     events inside the four STIMULUS BLOCKS

So each claim is now checked against the scope of the figure it sits beside.
"""
import glob
import os
import re

import numpy as np
import pandas as pd
from pptx import Presentation

DECK = (r"C:\Users\hsollim\Documents\HEAL_mini1p_SBI553\lab_meeting"
        r"\mini1p_SBI553_FINAL_2026-09-03_v2.pptx")
CSV = (r"C:\Users\hsollim\Documents\HEAL_mini1p_SBI553\lab_meeting"
       r"\figs_day2\Stats_per_mouse_allDeliveries.csv")
LOCO = (r"C:\Users\hsollim\Documents\HEAL_mini1p_SBI553\qc\locomotion"
        r"\Locomotion_summary.csv")
AFF = {1: "attending", 2: "lickbite", 3: "guarding", 4: "escape"}
REF = {1: "withdrawal", 2: "flinch"}
ORDER = ["withdrawal", "flinch", "attending", "lickbite", "guarding",
         "escape"]
W = {"withdrawal": 3.0, "flinch": 3.0, "attending": 10.0, "lickbite": 10.0,
     "guarding": 10.0, "escape": 10.0}
D = {"Day 1": r"C:\Users\hsollim\Documents\HEAL_mini1p_SBI553\videos"
                r"\output_corrected",
     "Day 2": r"C:\Users\hsollim\Documents\HEAL_mini1p_SBI553\videos"
                r"\output_day2_corrected"}


def s_(v, d=""):
    a = np.asarray(v).ravel()
    if a.size == 0:
        return d
    x = a[0]
    if isinstance(x, np.ndarray):
        x = x.ravel()[0] if x.size else d
    return str(x).strip()


def f_(v, d=np.nan):
    a = np.asarray(v, dtype=float).ravel()
    return float(a[0]) if a.size else d


def rates(folder, scope):
    out = {}
    for p in sorted(glob.glob(os.path.join(folder, "ScoringAB_*.mat"))):
        M = loadmat = __import__("scipy.io", fromlist=["loadmat"]).loadmat(p)
        fps = f_(M["frameRate"], 30.0)
        nU = int(f_(M.get("nUsed", 0), 0))
        sc = np.asarray(M["score"]).ravel().astype(int)
        sc = sc[:nU] if nU else sc
        N = len(sc)
        rx = np.asarray(M.get("reflexEvents", np.empty((0, 2)))).reshape(-1, 2)
        dF = np.asarray(M["dFrames"]).ravel().astype(int)
        dT = np.asarray(M["dTypes"]).ravel().astype(int)
        mid = s_(M.get("mouseID", "")).upper()
        ST = {}
        for c, nm in AFF.items():
            e = np.diff(np.concatenate(([0], (sc == c).astype(np.int8), [0])))
            ST[nm] = np.flatnonzero(e == 1) + 1
        for c, nm in REF.items():
            ff = rx[rx[:, 1] == c, 0].astype(int) if rx.size \
                else np.array([], int)
            ST[nm] = ff[(ff >= 1) & (ff <= N)]
        d = {}
        for b in ORDER:
            if scope == "session":
                tot, den = int(len(ST[b])), len(dF)
            else:
                tot, den = 0, 0
                for ty in range(1, 5):
                    sel = np.sort(dF[dT == ty])
                    if not len(sel):
                        continue
                    lo = int(sel[0])
                    hi = min(int(sel[-1] + W[b] * fps), N)
                    tot += int(np.sum((ST[b] >= lo) & (ST[b] < hi)))
                    den += len(sel)
            d[b] = tot / den if den else np.nan
        out[mid] = d
    return out


MED = {}
for scope in ("session", "blocks"):
    R1, R2 = rates(D["Day 1"], scope), rates(D["Day 2"], scope)
    mice = sorted(set(R1) & set(R2))
    MED[scope] = {b: float(np.median([R2[m][b] / R1[m][b] for m in mice
                                      if R1[m][b]])) for b in ORDER}

A = pd.read_csv(CSV)
A = A[A.stimulus == "ALL"]
nsig, nq = int((A.p_rate < .05).sum()), int((A.q_rate < .05).sum())
m2 = A[(A.mouse == "M2") & (A.behaviour == "escape")].iloc[0]
f1g = A[(A.mouse == "F1") & (A.behaviour == "guarding")].iloc[0]
dl = []
for folder in D.values():
    for p in sorted(glob.glob(os.path.join(folder, "ScoringAB_*.mat"))):
        from scipy.io import loadmat as _lm
        dl.append(len(np.asarray(_lm(p)["dFrames"]).ravel()))
L = pd.read_csv(LOCO)
bl = L[L.period == "baseline"].pivot_table(index="mouse", columns="day",
                                           values="speed_px_s")

CLAIMS = [
    ("slide1 lickbite median (session)", "0.20",
     f"{MED['session']['lickbite']:.2f}"),
    ("slide1 escape median (session)", "0.22",
     f"{MED['session']['escape']:.2f}"),
    ("slide1 withdrawal median (session)", "0.63",
     f"{MED['session']['withdrawal']:.2f}"),
    ("slide1 flinch median (session)", "0.82",
     f"{MED['session']['flinch']:.2f}"),
    ("slide1 attending median (session)", "0.51",
     f"{MED['session']['attending']:.2f}"),
    ("slide1 guarding median (session)", "0.45",
     f"{MED['session']['guarding']:.2f}"),
    ("tests significant", "25", str(nsig)),
    ("tests surviving BH", "22", str(nq)),
    ("M2 escape D1 events/stim", "60/83",
     f"{int(m2.n_events_day1)}/{int(m2.n_stim_day1)}"),
    ("M2 escape D2 events/stim", "13/70",
     f"{int(m2.n_events_day2)}/{int(m2.n_stim_day2)}"),
    ("M2 escape index (blocks)", "0.26", f"{m2.change_index:.2f}"),
    ("F1 guarding D1", "2/87",
     f"{int(f1g.n_events_day1)}/{int(f1g.n_stim_day1)}"),
    ("F1 guarding D2", "12/58",
     f"{int(f1g.n_events_day2)}/{int(f1g.n_stim_day2)}"),
    ("stimuli range", "57-100", f"{min(dl)}-{max(dl)}"),
    ("baseline speed ratio", "0.47",
     f"{bl.iloc[:, 1].mean() / bl.iloc[:, 0].mean():.2f}"),
    ("baseline mice down", "6/6",
     f"{int((bl.iloc[:, 1] < bl.iloc[:, 0]).sum())}/6"),
]

print("=== claim vs source, scope-aware ===")
bad = 0
for name, in_deck, truth in CLAIMS:
    ok = in_deck == truth
    bad += not ok
    print(f"  {'OK  ' if ok else 'FAIL'} {name:36s} deck {in_deck:>8s}   "
          f"source {truth:>8s}")

print("\n=== does the deck state each figure's scope? ===")
prs = Presentation(DECK)
txt = {i: " ".join(sh.text_frame.text for sh in sl.shapes
                   if sh.has_text_frame)
       for i, sl in enumerate(prs.slides, 1)}
need = [(2, "WHOLE SESSION"), (5, "STIMULUS BLOCKS")]
for i, phrase in need:
    ok = phrase.lower() in txt.get(i, "").lower()
    bad += not ok
    print(f"  {'OK  ' if ok else 'FAIL'} slide {i} states '{phrase}'")

print("\n=== jargon / length check on titles ===")
JARG = ["dissociation", "biased agonist", "supraspinal", "spinally"]
for i, t in txt.items():
    first = t.split("\n")[0] if t else ""
    hits = [j for j in JARG if j in first.lower()]
    if hits:
        print(f"  note slide {i} title contains {hits}")
    if len(first) > 62 and i > 1:
        print(f"  note slide {i} title is {len(first)} chars: {first[:70]}")

print("\n" + "=" * 66)
print("AUDIT PASSED - every deck claim matches the source"
      if not bad else f"AUDIT FAILED: {bad} problem(s)")
