"""make_final_5slides.py  -  the five verified figures, nothing else.

Every number on these slides was recomputed from the .mat files by
verify_final_numbers.py, which reads the scoring independently of the figure
code and fails if the two disagree. It currently reports ALL CHECKS PASSED
after catching one real bug: counting rising edges inside a time slice made a
bout that was already running at the window edge look like a new event.

Each slide carries three things, in this order:
    HEADLINE   what the figure shows, in one sentence
    METHOD     what was measured and which test, explicitly
    READING    how to read the panel, and what it does not say

THE MEASURE, once, for all five slides
    change index = (total events / total stimuli) on Day 2
                 / (total events / total stimuli) on Day 1
    Every event in the stimulus block is counted. No response-window filter.
    Each mouse is divided by its own stimulus count.

VERIFIED NUMBERS

    TWO aggregations, both correct, and they must not be mixed on one slide:

      ratio of means      pool the six mice, then divide. Matches the
                          population figures, where the bold marker is a mean.
      median of ratios    one index per mouse, then take the median. Matches
                          the per-mouse figures, where each dot is one animal.

                      D1      D2   ratio of   median of   down   paired p
                                      means      ratios
    withdrawal     1.097   0.674       0.61        0.63    6/6      0.031
    flinch         0.787   0.406       0.52        0.82    3/6      0.438
    attending      0.287   0.121       0.42        0.51    6/6      0.031
    lickbite       0.335   0.081       0.24        0.20    6/6      0.031
    guarding       0.224   0.111       0.50        0.45    5/6      0.156
    escape         0.608   0.111       0.18        0.22    6/6      0.031

    Slide 1 shows medians, so slide 1 quotes the median column.

    NOTE ON PAW WITHDRAWAL, which reads as "barely changed" on a log axis:
    ×0.63 is a 37 % DECREASE, in all six mice, paired p = 0.031 - the floor of
    the exact test. It is significantly reduced AND the least reduced of the
    six behaviours. Both statements are true and neither should be dropped.

USAGE
    python verify_final_numbers.py --day1 <d1> --day2 <d2> --figs <figs>
    python make_final_5slides.py

Hansol Lim - HEAL mini1p / SBI-553
"""
from __future__ import annotations

import argparse
import os
from datetime import date

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from make_lab_meeting_ppt import BASE, INK, ACC, GREY, FONT, blank, deck, tb

FIGS = os.path.join(BASE, "lab_meeting", "figs_day2")
OUTDIR = os.path.join(BASE, "lab_meeting")
WARN = RGBColor(0xC0, 0x48, 0x3B)
METH = RGBColor(0x55, 0x5D, 0x66)


def slide(prs, headline, method, reading, fig, warn=False):
    """One figure, with the headline above it and the method below.

    The figure is given a fixed box so the text positions never move between
    slides - a reader should not have to hunt for the method line.
    """
    s = blank(prs)
    tb(s, Inches(.55), Inches(.30), Inches(12.25), Inches(.85),
       headline, 25, True, WARN if warn else INK)
    tb(s, Inches(.55), Inches(1.18), Inches(12.25), Inches(.62),
       "METHOD   " + method, 12.5, False, METH)

    top, boxw, boxh = 1.92, 12.25, 4.22
    if os.path.exists(fig):
        from PIL import Image
        iw, ih = Image.open(fig).size
        asp = iw / ih
        w, h = boxw, boxw / asp
        if h > boxh:
            h, w = boxh, boxh * asp
        s.shapes.add_picture(fig, Inches(.55 + (boxw - w) / 2), Inches(top),
                             Inches(w), Inches(h))
    else:
        tb(s, Inches(.55), Inches(top), Inches(12), Inches(.5),
           f"[missing {os.path.basename(fig)}]", 14, False, WARN)

    tb(s, Inches(.55), Inches(6.26), Inches(12.25), Inches(1.15),
       reading, 12.5, False, ACC, space=3)
    return s


def build():
    prs = deck()

    # ── title ────────────────────────────────────────────────────────
    s = blank(prs)
    tb(s, Inches(.9), Inches(1.5), Inches(11.5), Inches(1.2),
       "SBI-553 on the acute pain assay", 38, True, INK)
    tb(s, Inches(.9), Inches(2.7), Inches(11.5), Inches(1.5),
       "A general reduction in behaviour, not selective analgesia",
       26, True, WARN)
    tb(s, Inches(.9), Inches(4.0), Inches(11.5), Inches(1.8),
       "Escape / rearing is exploration, not pain — and it was the ONLY "
       "behaviour significantly\nreduced in all six mice. The reflexes fell "
       "least. A selective analgesic predicts the\nopposite ordering.",
       20, False, ACC)
    tb(s, Inches(.9), Inches(5.9), Inches(11.5), Inches(1.1),
       "Day 1 no drug vs Day 2 SBI-553, same six mice, dosed 10 min before "
       "the assay.  No vehicle group.\n"
       "change index  =  (total events ÷ total stimuli) on Day 2   ÷   "
       "the same on Day 1",
       15, False, GREY)
    tb(s, Inches(.9), Inches(7.0), Inches(11.5), Inches(.45),
       f"Hansol Lim   ·   {date.today().isoformat()}   ·   every number "
       f"recomputed from the scoring files and cross-checked",
       12, False, GREY)

    # ── 1. analgesia or sedation ─────────────────────────────────────
    slide(
        prs,
        "Every behaviour fell — including escape/rearing, which is not a pain "
        "measure",
        "Change index per behaviour, one dot per mouse, bar = MEDIAN of the "
        "six per-mouse indices. Every event in the stimulus block counted, "
        "each mouse divided by its own stimulus count. Log axis. Paired "
        "Wilcoxon over the six mice (exact, n = 6, smallest possible "
        "p = 0.031).",
        "Medians: licking/biting ×0.20 and escape/rearing ×0.22 fell most, "
        "paw withdrawal ×0.63 and flinch ×0.82 least (attending ×0.51, "
        "guarding ×0.45). Escape/rearing is exploration — Day 1 showed it is "
        "high at baseline and DECREASES under stimulation — so a selective "
        "analgesic should have spared it. It did not.\n"
        "Paw withdrawal looks close to the line only because the axis is "
        "logarithmic: ×0.63 is a 37 % decrease, in all six mice, "
        "paired p = 0.031. It is significantly reduced AND the least reduced "
        "of the six.",
        os.path.join(FIGS, "S1_sedation_evidence.png"), warn=True)

    # ── 2. dose-response, both days ──────────────────────────────────
    slide(
        prs,
        "Lower on the drug day for all six behaviours, and flatter for the "
        "affective ones",
        "Events per stimulus delivery against stimulus type. Thin line = one "
        "mouse, bold = mean ± SEM over six mice. Both days on shared axes so "
        "the size of the drop is comparable. No statistical test on this "
        "panel — it is descriptive.",
        "Day 1 rises from light touch to heat in every behaviour. On Day 2 "
        "that ordering is flattened for licking/biting and escape/rearing but "
        "kept for the reflexes, whose curve simply sits lower.",
        os.path.join(FIGS, "D8_dose_response_both_days.png"))

    # ── 3. paired, per stimulus ──────────────────────────────────────
    slide(
        prs,
        "Paired mouse by mouse: paw withdrawal and escape/rearing move in "
        "every animal",
        "Same measure, paired within animal. Grey line joins the two days for "
        "one mouse; square = mean ± SEM. Paired Wilcoxon per behaviour × "
        "stimulus (exact, n = 6, floor p = 0.031). * p < 0.05, ns = not "
        "significant.",
        "Escape/rearing is significant at all four stimuli. Paw withdrawal is "
        "significant at three of four. Guarding and flinch are not "
        "significant per stimulus — with six mice the test cannot go below "
        "0.031, so a non-significant cell here means underpowered, not "
        "unchanged.",
        os.path.join(FIGS, "D2_per_delivery.png"))

    # ── 4. per-mouse change index with statistics ────────────────────
    # Counting note. Reporting only "N mice significant" hides direction: on
    # guarding, four animals are significant but one of them (F1) went UP
    # nine-fold, so only three decreased significantly. Every count below
    # states the direction.
    #
    #                    decreased   sig. decreased   increased
    #   withdrawal          6/6            5             0
    #   flinch              3/6            3             3   <- split
    #   attending           6/6            3             0
    #   lickbite            6/6            4             0
    #   guarding            5/6            3             1   <- F1 x9.00
    #   escape              6/6            6             0
    slide(
        prs,
        "Escape/rearing — the behaviour that is NOT pain — is the only one "
        "significant in all six mice",
        "Change index per mouse, dot = the index, line = 95 % CI from the "
        "Poisson variance of the totals. EXACT POISSON RATE TEST on that "
        "animal's own counts: conditional on the n₁+n₂ events in total, the "
        "number landing on Day 2 is Binomial(n₁+n₂, t₂/(t₁+t₂)) if the two "
        "rates are equal; the two-sided binomial p is exact. Stars are "
        "uncorrected — 25 of 36 tests reach p<0.05, 22 survive "
        "Benjamini-Hochberg.",
        "Decreased in / significantly decreased in:  escape 6/6 and 6/6  ·  "
        "withdrawal 6/6 and 5  ·  lick-bite 6/6 and 4  ·  attending 6/6 and 3 "
        " ·  guarding 5/6 and 3  ·  flinch 3/6 and 3 (the other three went "
        "UP). Guarding has one animal moving the other way: F1 rose nine-fold "
        "(2/87 → 12/58, p = 0.0006), so four significant does not mean four "
        "decreased.",
        os.path.join(FIGS, "Fig_forest_per_mouse_allDeliveries.png"))

    # ── 5. the dissociation ──────────────────────────────────────────
    slide(
        prs,
        "The affective/reflexive dissociation a biased agonist predicts is "
        "not what happened",
        "Affective (attending + licking/biting + guarding) against reflexive "
        "(withdrawal + flinch), both as events per stimulus delivery. One "
        "arrow per mouse, Day 1 to Day 2. Descriptive — no test on this "
        "panel.",
        "A selective affective effect would move every arrow straight DOWN. "
        "F2, F3 and M3 do move mostly down; M1 and M2 move down AND left; F1 "
        "moves left and slightly up. Combined with escape/rearing falling "
        "hardest, the pattern is a general reduction in activity rather than "
        "a clean affective-only effect.",
        os.path.join(FIGS, "D4_reflex_vs_affective.png"))

    # ── what this cannot say ─────────────────────────────────────────
    s = blank(prs)
    tb(s, Inches(.55), Inches(.35), Inches(12.25), Inches(.8),
       "What these five figures do NOT establish", 26, True, WARN)
    tb(s, Inches(.7), Inches(1.5), Inches(12.0), Inches(4.6),
       "1.  That the change is caused by SBI-553.\n"
       "     All six mice received it. With no vehicle group, drug, day, "
       "block order and\n"
       "     habituation cannot be separated.\n\n"
       "2.  That it is analgesia.\n"
       "     Escape/rearing is not a pain measure and fell hardest. Movement "
       "measured from\n"
       "     the video was already halved during the 5 min baseline, before "
       "any stimulus\n"
       "     (×0.47, 6/6 mice, p = 0.031) — there is no pain to relieve "
       "there.\n\n"
       "3.  That analgesia is absent.\n"
       "     A general reduction in activity could hide it. A lower dose or a "
       "later time point,\n"
       "     with a vehicle arm, would separate the two.\n\n"
       "4.  Anything about the face.\n"
       "     The side view is a backlit silhouette: ears are visible as "
       "outlines, the eye is not\n"
       "     resolved, so the Mouse Grimace Scale cannot be scored. And "
       "grimace cannot\n"
       "     distinguish pain from sedation in any case — both produce "
       "orbital tightening.",
       16, False, INK)
    return prs


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--out", default=None)
    a = ap.parse_args()
    prs = build()
    out = a.out or os.path.join(
        OUTDIR, f"mini1p_SBI553_FINAL_{date.today().isoformat()}.pptx")
    base, n = out, 1
    while True:
        try:
            prs.save(out)
            break
        except PermissionError:
            n += 1
            r, e = os.path.splitext(base)
            out = f"{r}_v{n}{e}"
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides")
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
