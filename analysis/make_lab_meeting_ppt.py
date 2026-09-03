"""make_lab_meeting_ppt.py  -  build the lab-meeting deck.

Simple on purpose: one message per slide, two colours, one font, no clip art.
The embedded video is the clip you scored yourself.

USAGE
    python make_lab_meeting_ppt.py
    python make_lab_meeting_ppt.py --day2 <day2 corrected folder>

OUTPUT
    lab_meeting\\mini1p_SBI553_labmeeting_<date>.pptx

Hansol Lim - HEAL mini1p / SBI-553
"""
from __future__ import annotations

import argparse
import glob
import os
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = r"C:\Users\hsollim\Documents\HEAL_mini1p_SBI553"
FIGS = os.path.join(BASE, "lab_meeting", "figs")
CLIPS = os.path.join(BASE, "ppt_clips")
OUTDIR = os.path.join(BASE, "lab_meeting")

INK = RGBColor(0x1F, 0x29, 0x37)
ACC = RGBColor(0x1C, 0x6E, 0x8C)
GREY = RGBColor(0x6B, 0x72, 0x80)
FONT = "Calibri"
W, H = Inches(13.333), Inches(7.5)


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def tb(slide, x, y, w, h, text, size=18, bold=False, color=INK,
       align=PP_ALIGN.LEFT, space=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space)
        r = para.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return box


def title_slide(prs, sub):
    s = blank(prs)
    tb(s, Inches(.9), Inches(2.2), Inches(11.5), Inches(1.4),
       "Acute pain assay under mini1p imaging", 40, True, INK)
    tb(s, Inches(.9), Inches(3.5), Inches(11.5), Inches(1.2),
       "Day 1 baseline, six mice  \u2014  what the behaviour looks like\n"
       "before we add SBI-553", 22, False, ACC)
    tb(s, Inches(.9), Inches(5.6), Inches(11.5), Inches(1.0),
       f"Hansol Lim   \u00b7   {sub}", 15, False, GREY)
    return s


def head(slide, title, subtitle=None):
    tb(slide, Inches(.6), Inches(.35), Inches(12.1), Inches(.7),
       title, 28, True, INK)
    if subtitle:
        tb(slide, Inches(.6), Inches(1.05), Inches(12.1), Inches(.6),
           subtitle, 15, False, GREY)


def fig_slide(prs, title, fig, subtitle=None, note=None, top=1.55):
    s = blank(prs)
    head(s, title, subtitle)
    if os.path.exists(fig):
        # fit inside the content box, preserving aspect
        from PIL import Image
        iw, ih = Image.open(fig).size
        boxw = 12.4
        boxh = (7.2 - top) if note is None else (6.55 - top)
        # Fit by aspect ratio only. Converting pixels through an assumed DPI
        # is fragile because bbox_inches="tight" changes the saved size, and
        # it shrank every figure to about half the slide.
        asp = iw / ih
        w, h = boxw, boxw / asp
        if h > boxh:
            h, w = boxh, boxh * asp
        s.shapes.add_picture(fig, Inches(.45 + (boxw - w) / 2), Inches(top),
                             Inches(w), Inches(h))
    else:
        tb(s, Inches(.6), Inches(top), Inches(12.1), Inches(1),
           f"[missing: {os.path.basename(fig)}]", 16, False, GREY)
    if note:
        tb(s, Inches(.6), Inches(6.62), Inches(12.1), Inches(.7),
           note, 14, False, ACC)
    return s


def bullets_slide(prs, title, rows, subtitle=None):
    s = blank(prs)
    head(s, title, subtitle)
    y = 1.9
    for txt, sub in rows:
        tb(s, Inches(.8), Inches(y), Inches(11.7), Inches(.45), txt, 20, True, INK)
        if sub:
            tb(s, Inches(1.1), Inches(y + .42), Inches(11.4), Inches(.5),
               sub, 15, False, GREY)
            y += 1.02
        else:
            y += .62
    return s


def table_slide(prs, title, header, rows, subtitle=None, colw=None):
    s = blank(prs)
    head(s, title, subtitle)
    nr, nc = len(rows) + 1, len(header)
    tw = Inches(12.1)
    shp = s.shapes.add_table(nr, nc, Inches(.6), Inches(1.85), tw,
                             Inches(.4 * nr)).table
    if colw:
        tot = sum(colw)
        for i, cw in enumerate(colw):
            shp.columns[i].width = Inches(12.1 * cw / tot)
    def put(c, txt, size, bold):
        # add_run rather than cell.text: an empty string leaves no run at all,
        # and then styling it raises IndexError
        pr = c.text_frame.paragraphs[0]
        r = pr.add_run()
        r.text = str(txt)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = INK

    for j, htxt in enumerate(header):
        put(shp.cell(0, j), htxt, 14, True)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            put(shp.cell(i, j), v, 13, False)
    return s


def video_slide(prs, title, mp4, poster, subtitle=None, note=None):
    s = blank(prs)
    head(s, title, subtitle)
    if mp4 and os.path.exists(mp4):
        from PIL import Image
        iw, ih = Image.open(poster).size if poster and os.path.exists(poster) \
            else (1206, 730)
        boxw, boxh = 11.8, 4.55
        asp = iw / ih
        w, h = boxw, boxw / asp
        if h > boxh:
            h, w = boxh, boxh * asp
        s.shapes.add_movie(mp4, Inches(.45 + (boxw - w) / 2), Inches(1.75),
                           Inches(w), Inches(h),
                           poster_frame_image=poster if poster and
                           os.path.exists(poster) else None)
        tb(s, Inches(.6), Inches(6.5), Inches(12.1), Inches(.6),
           note or "click to play", 14, False, ACC)
    else:
        tb(s, Inches(.6), Inches(2.0), Inches(12.1), Inches(1),
           "[video not found]", 16, False, GREY)
    return s


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--day2", default=None)
    ap.add_argument("--out", default=None)
    a = ap.parse_args()
    have2 = bool(a.day2 and glob.glob(os.path.join(a.day2, "ScoringAB_*.mat")))
    os.makedirs(OUTDIR, exist_ok=True)

    mp4 = sorted(glob.glob(os.path.join(CLIPS, "PPT_*.mp4")))
    mp4 = mp4[0] if mp4 else None
    poster = None
    if mp4:
        cands = sorted(glob.glob(mp4[:-4] + "_frame*.png"))
        poster = cands[len(cands) // 2] if cands else None

    prs = deck()
    today = date.today().isoformat()

    # 1 title
    title_slide(prs, f"lab meeting, {today}")

    # 2 the question
    bullets_slide(
        prs, "What we are building",
        [("A behavioural read-out of acute pain that we can pair with mini1p imaging",
          "reflexive and affective-motivational responses scored separately \u2014 "
          "never combined into one pain score"),
         ("Four stimuli spanning innocuous to noxious, in one session per mouse",
          "Light touch \u00b7 Mild touch \u00b7 Heat \u00b7 Pin prick, order randomised"),
         ("Then ask whether SBI-553 changes it",
          "Day 1 no drug (done) \u2192 Day 2 drug, same six mice, within-animal")],
        "the assay first, the drug second")

    # 3 design
    fig_slide(prs, "Session design",
              os.path.join(FIGS, "F1_design.png"),
              "one 28 min session per mouse",
              "Baseline first, so every stimulus block has its own within-session "
              "reference. Blocks run on the clock: 5\u201310, 11\u201316, 17\u201322, 23\u201328 min.")

    # 4 how we score - video
    video_slide(prs, "How it is scored",
                mp4, poster,
                "both cameras at once, six behaviours, live label track",
                "Bottom camera shows the applicator, side camera shows the "
                "behaviour. Black lines = stimulus delivered. "
                "Dots = brief marks (< 0.25 s).")

    # 5 the six behaviours
    table_slide(
        prs, "Six behaviours, two classes",
        ["", "behaviour", "what it is"],
        [["reflexive", "Paw withdrawal", "rapid lift of the stimulated paw at contact"],
         ["", "Flinch", "brief shake or flick at contact"],
         ["affective", "Paw attending", "orients to and inspects the paw, no mouth contact"],
         ["", "Licking / biting", "any mouth contact with the paw"],
         ["", "Guarding", "paw held up, no weight-bearing"],
         ["", "Escape / rearing", "locomotor escape, or rearing"]],
        "reported separately \u2014 a single \u201cpain score\u201d would hide the distinction",
        colw=[1.4, 3.0, 7.7])

    # 6 result: dose response
    fig_slide(prs, "Day 1: the assay separates the stimuli",
              os.path.join(FIGS, "F2_dose_response.png"),
              "baseline-subtracted, 6 mice, median + individual animals",
              "Light touch gives the smallest response and Heat the largest, for "
              "every behaviour. Pin prick sits between Mild touch and Heat.")

    # 7 per delivery
    fig_slide(prs, "Same result, per stimulus delivered",
              os.path.join(FIGS, "F3_per_delivery.png"),
              "controls for the fact that blocks did not all get the same number "
              "of deliveries (16\u201331)",
              "Heat evokes a withdrawal on ~1.5 of every 2 deliveries; "
              "light touch on ~0.7.")

    # 6b raw counts
    fig_slide(prs, "The same thing as raw counts",
              os.path.join(FIGS, "F2b_event_counts.png"),
              "how many events happened in each 5 min block",
              "Every block is a fixed 300 s, so these are directly comparable "
              "with no normalisation at all. Heat: ~28 withdrawals and ~13 "
              "licking bouts per block; light touch ~16 and ~2.")

    # 7b peri-stimulus time histogram
    fig_slide(prs, "Reflexive and affective separate in TIME",
              os.path.join(FIGS, "E1_psth.png"),
              "every delivery aligned at t = 0, probability the behaviour is "
              "occurring",
              "Withdrawal and flinch spike exactly at contact. Licking and "
              "guarding build over 1-8 s. This is the same alignment we will "
              "use for peri-event dF/F, so behaviour and imaging can go side "
              "by side.")

    # 7c response probability
    fig_slide(prs, "Response probability per delivery",
              os.path.join(FIGS, "E2_response_probability.png"),
              "did the behaviour occur within 3 s of the stimulus?",
              "Heat evokes a withdrawal on essentially every delivery "
              "(P = 1.00) against 0.74 for light touch. This is one trial per "
              "delivery - 3,204 trials instead of 24 block means, which is "
              "where the statistical power actually is.")

    # 7d response probability, 10 s on isolated deliveries
    fig_slide(prs, "Same measure with a 10 s window",
              os.path.join(FIGS, "E2b_response_probability_10s.png"),
              "only deliveries with 10 s of clearance before the next one "
              "(122 of 535)",
              "3 s catches the reflexes but cuts the slower affective response "
              "short; 10 s catches it but only 23 % of deliveries have 10 s of "
              "clear space, because the median gap between deliveries is 2.7 s.")

    # 8 escape is exploration
    fig_slide(prs, "One behaviour runs the other way",
              os.path.join(FIGS, "F4_baseline_vs_block.png"),
              "escape / rearing is the only behaviour with a high baseline",
              "It DECREASES during stimulation. That reads as exploration being "
              "suppressed, not as a pain response \u2014 so it should not be pooled "
              "with the others.")

    # 9 per mouse
    fig_slide(prs, "Is it consistent across animals?",
              os.path.join(FIGS, "F5_per_mouse.png"),
              "one line per mouse, all six shown",
              "The ordering holds in most animals. With n = 6 the spread is what "
              "it is \u2014 individual variation is the honest picture, not noise to "
              "average away.")

    # 10 day 1 vs day 2
    fig_slide(prs, "Day 1 vs Day 2 (drug)" if have2
              else "Ready for Day 2",
              os.path.join(FIGS, "F6_day1_vs_day2_counts.png"),
              "whole-session event rate, one line per mouse"
              if have2 else "Day 1 distribution; the drug day drops into the "
                            "same figure",
              "Primary contrast is within-animal: each mouse is its own control."
              if have2 else
              "Analysis is written and validated already \u2014 tested against a "
              "synthetic drug day, it recovered the injected effect and reported "
              "no change elsewhere.")

    # 10b within-block time course
    fig_slide(prs, "Does the response habituate within a block?",
              os.path.join(FIGS, "E3_within_block_timecourse.png"),
              "each 5 min block split into 1 min bins",
              "Worth watching on the drug day: SBI-553 could change the SHAPE "
              "of the response rather than its level.")

    # 10c block position control
    fig_slide(prs, "Control: is it the stimulus, or the time in the session?",
              os.path.join(FIGS, "E4_block_position.png"),
              "response by block position (1st to 4th) instead of by stimulus",
              "Stimulus order was randomised per mouse, so if the response "
              "tracks the stimulus and not the position, the randomisation did "
              "its job.")

    # 10d reflexive vs affective
    fig_slide(prs, "Do the two classes dissociate?",
              os.path.join(FIGS, "E5_reflex_vs_affective.png"),
              "one point per mouse, whole-session totals",
              "If reflexive and affective totals were interchangeable there "
              "would be no point scoring them separately.")

    # 11 honest limits
    bullets_slide(
        prs, "What to keep in mind",
        [("Counts and rates, not durations",
          "the sessions were scored by tapping, so a mark records the keypress, "
          "not how long the behaviour lasted. Guarding carries a nominal 1 s."),
         ("n = 6 sets a hard floor on p",
          "the exact paired test cannot go below p = 0.031, so one pre-specified "
          "outcome can reach significance and a whole table cannot"),
         ("Escape / rearing is not a pain measure here",
          "high baseline, decreases under stimulation \u2014 report it separately"),
         ("The notch at t = 0 in the peri-stimulus plots is us, not the mouse",
          "to tap a delivery key you have to let go of the behaviour key"),
         ("Every correction is logged",
          "93 changes across the six sessions, each one recorded; the raw "
          "scoring is never modified")],
        "so nobody is surprised later")

    # 12 next
    bullets_slide(
        prs, "Next",
        [("Day 2 with SBI-553, same six mice", "figures update automatically"),
         ("One piece of hardware would fix three problems at once",
          "a TTL or LED on the applicator gives camera sync, frame-accurate "
          "delivery times, and stimulus identity \u2014 which image analysis cannot "
          "recover, the four applicators look identical from below"),
         ("A second baseline at the end of the session",
          "5 min, and it makes drift over the 28 min estimable"),
         ("Automatic scoring: not yet",
          "rearing is detectable (d\u2032 = 6.0) but not calibrated; the others "
          "do not have enough labelled events")])

    out = a.out or os.path.join(
        OUTDIR, f"mini1p_SBI553_labmeeting_{today}.pptx")
    prs.save(out)
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides")
    print(f"wrote {out}")
    if mp4:
        print(f"  embedded video: {os.path.basename(mp4)}")
    else:
        print("  NOTE: no PPT_*.mp4 found in ppt_clips - the video slide is empty")


if __name__ == "__main__":
    main()
