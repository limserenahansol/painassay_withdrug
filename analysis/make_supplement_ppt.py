"""make_supplement_ppt.py  -  a separate deck: how to read the numbers.

WHY A SEPARATE FILE
    The main deck stays untouched so it can be open and edited while this is
    generated. This one covers three things the main deck assumes:

      1  how the rate ratio is calculated, worked through on real numbers
      2  how to read the per-mouse forest plot
      3  Day 2 on its own, plus the video-measured activity result

USAGE
    python make_supplement_ppt.py
    python make_supplement_ppt.py --out <path.pptx>

Hansol Lim - HEAL mini1p / SBI-553
"""
from __future__ import annotations

import argparse
import os
from datetime import date

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from make_lab_meeting_ppt import (BASE, INK, ACC, GREY, FONT, blank,
                                  bullets_slide, deck, fig_slide, head,
                                  table_slide, tb)

D2 = os.path.join(BASE, "lab_meeting", "figs_day2")
D2ONLY = os.path.join(BASE, "lab_meeting", "figs_day2only")
OUTDIR = os.path.join(BASE, "lab_meeting")
WARN = RGBColor(0xC0, 0x48, 0x3B)
OK = RGBColor(0x1C, 0x6E, 0x8C)
MONO = "Consolas"


def mono(slide, x, y, w, h, text, size=15, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(2)
        r = para.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.name = MONO
        r.font.color.rgb = color
    return box


def divider(prs, title, sub):
    s = blank(prs)
    tb(s, Inches(.9), Inches(2.7), Inches(11.5), Inches(1.2), title, 34,
       True, INK)
    tb(s, Inches(.9), Inches(4.0), Inches(11.5), Inches(1.4), sub, 19, False,
       ACC)
    return s


def build():
    prs = deck()

    s = blank(prs)
    tb(s, Inches(.9), Inches(2.3), Inches(11.5), Inches(1.3),
       "How to read the Day 2 numbers", 38, True, INK)
    tb(s, Inches(.9), Inches(3.6), Inches(11.5), Inches(1.2),
       "Supplement to the main deck  ·  the rate ratio, the forest plot, "
       "Day 2 on its own,\nand the activity measured from the video",
       20, False, ACC)
    tb(s, Inches(.9), Inches(5.8), Inches(11.5), Inches(.8),
       f"Hansol Lim   ·   {date.today().isoformat()}", 15, False, GREY)

    # ─────────── 1. the rate ratio, worked ──────────────────────────────
    divider(prs, "1.  The rate ratio",
            "one number per mouse per behaviour, and how it is built")

    s = blank(prs)
    head(s, "Step 1: count events, and count the stimuli that caused them",
         "real numbers, mouse M2, escape / rearing")
    mono(s, Inches(.8), Inches(1.9), Inches(11.7), Inches(2.2),
         "                       Day 1 no drug     Day 2 SBI-553\n"
         "\n"
         "  escape events              69                6\n"
         "  stimuli delivered          83               70\n",
         17)
    tb(s, Inches(.8), Inches(4.2), Inches(11.7), Inches(1.8),
       "The stimulus is delivered by hand, so the number of taps is NOT the "
       "same on the two days:\n83 on Day 1, 70 on Day 2. Comparing 69 events "
       "against 6 events would mix the change in\nbehaviour with the change "
       "in how many chances the animal was given.",
       17, False, INK)

    s = blank(prs)
    head(s, "Step 2: divide by that animal's own stimulus count",
         "this is the number every figure uses")
    mono(s, Inches(.8), Inches(1.9), Inches(11.7), Inches(2.6),
         "  Day 1     69 events  /  83 stimuli   =   0.83  events per stimulus\n"
         "  Day 2      6 events  /  70 stimuli   =   0.09  events per stimulus\n",
         18)
    tb(s, Inches(.8), Inches(3.6), Inches(11.7), Inches(1.0),
       "In words: on Day 1 this mouse tried to escape about 8 times out of "
       "every 10 stimuli.\nOn Day 2, about 1 time out of every 10.",
       17, False, OK)
    tb(s, Inches(.8), Inches(5.0), Inches(11.7), Inches(1.4),
       "This is what the axis label “total event number / total stimulus "
       "delivery” means.\nEach mouse is divided by ITS OWN count, never "
       "by a group average.",
       15, False, GREY)

    s = blank(prs)
    head(s, "Step 3: the rate ratio is Day 2 divided by Day 1",
         "so 1 means no change")
    mono(s, Inches(.8), Inches(1.9), Inches(11.7), Inches(1.6),
         "  rate ratio  =  0.09  /  0.83   =   0.10\n",
         20)
    table_slide_rows = None
    tb(s, Inches(.8), Inches(3.1), Inches(11.7), Inches(2.6),
       "Read it as a multiplier:\n\n"
       "    1.0     no change\n"
       "    0.5     halved\n"
       "    0.10    down to a tenth  —  this mouse\n"
       "    2.0     doubled\n",
       17, False, INK)
    tb(s, Inches(.8), Inches(5.9), Inches(11.7), Inches(1.0),
       "Why a ratio and not a per-cent change: per cent is bounded at "
       "−100 % but unbounded upwards,\nso one animal that went up 650 % "
       "flattens every other bar. A ratio on a log axis treats halving\n"
       "and doubling as the same distance.",
       15, False, GREY)

    table_slide(
        prs, "Why the per-mouse test is possible at all",
        ["", "value", "what it buys"],
        [["Stimuli per mouse per day", "58 to 100",
          "each delivery is one trial, so a single animal has ~60-100 trials"],
         ["Test used", "Mann-Whitney on that mouse's own deliveries",
          "no group needed - the mouse is its own experiment"],
         ["Confidence interval", "from the Poisson variance of the totals",
          "the horizontal line in the forest plot"],
         ["If we had used block means instead", "4 numbers per mouse per day",
          "no per-mouse test would be possible at all"]],
        "this is the payoff of scoring every delivery rather than every block",
        colw=[3.0, 3.4, 5.6])

    # ─────────── 2. the forest plot ─────────────────────────────────────
    divider(prs, "2.  The forest plot",
            "one row per mouse, one panel per behaviour")

    s = blank(prs)
    head(s, "How to read it")
    rows = [
        ("Each row is one mouse", "F1 at the top, M3 at the bottom, the SAME "
                                  "order in every panel - so you can follow "
                                  "one animal left to right"),
        ("The dot is that mouse's rate ratio",
         "Day 2 ÷ Day 1 in events per stimulus"),
        ("The horizontal line is the 95 % confidence interval",
         "a short line means that mouse gave many consistent trials; a long "
         "line means few or scattered"),
        ("The dashed vertical line is 1 = no change",
         "left of it = fewer events on the drug day"),
        ("Red means p < 0.05 for THAT MOUSE ALONE",
         "not the group - this is a within-animal test on its own ~60-100 "
         "deliveries"),
        ("The axis is logarithmic",
         "1/4 is as far from 1 as 4 is, so increases and decreases are "
         "comparable"),
    ]
    y = 1.65
    for t, sub in rows:
        tb(s, Inches(.8), Inches(y), Inches(11.7), Inches(.4), t, 17, True,
           INK)
        tb(s, Inches(1.1), Inches(y + .38), Inches(11.4), Inches(.45), sub,
           14, False, GREY)
        y += .92

    fig_slide(prs, "The forest plot itself",
              os.path.join(D2, "Fig_forest_per_mouse_allDeliveries.png"),
              None,
              "Read across a row to follow one animal. M1 and M2 are red in "
              "every panel. F1 is the exception: unchanged or up for the pain "
              "behaviours, but still down for escape / rearing.",
              top=1.15)
    fig_slide(prs, "The same thing without confidence intervals",
              os.path.join(D2, "D5_per_mouse_change.png"), None,
              "Simpler version of the previous slide: dot and stick, no "
              "statistics. Useful when the point is only the direction and "
              "size.", top=1.15)

    # ─────────── 3. Day 2 on its own ────────────────────────────────────
    divider(prs, "3.  Day 2 on its own",
            "the same descriptive figures we made for Day 1, "
            "computed from the drug day only")
    for name, title, note in (
        ("F2_dose_response.png", "Day 2 only: response by stimulus",
         "the Day-1 ordering (light touch smallest, heat largest) is largely "
         "flattened"),
        ("F3_per_delivery.png",
         "Day 2 only: total event number / total stimulus delivery",
         "the comparable version - each mouse divided by its own delivery "
         "count"),
        ("F4_baseline_vs_block.png",
         "Day 2 only: stimulus blocks against each animal's own baseline",
         "worth comparing with the Day-1 version: escape / rearing no longer "
         "has a high baseline to fall from"),
        ("F5_per_mouse.png", "Day 2 only: every mouse individually",
         "one panel per animal, no averaging"),
    ):
        p = os.path.join(D2ONLY, name)
        if os.path.exists(p):
            fig_slide(prs, title, p, None, note, top=1.15)

    # ─────────── 4. the video measure ───────────────────────────────────
    divider(prs, "4.  Activity measured from the video",
            "independent of the manual scoring - it comes from the pixels")
    fig_slide(prs, "The baseline is the decisive comparison",
              os.path.join(D2, "V1_locomotion.png"), None,
              "No stimulus is delivered during the baseline, so there is no "
              "pain to relieve. Movement was already halved there in 6/6 "
              "mice (p = 0.031). That cannot be analgesia.",
              top=1.15)
    table_slide(
        prs, "What the video measured, baseline only",
        ["measure", "Day 1", "Day 2", "ratio", "mice down", "p"],
        [["centroid speed (px/s)", "31.7", "14.9", "0.47", "6 / 6", "0.031"],
         ["fraction of time moving", "0.72", "0.52", "0.73", "6 / 6", "0.031"],
         ["silhouette area (posture)", "30536", "23167", "0.76", "6 / 6",
          "0.031"]],
        "the first 5 minutes of the session, before any stimulus",
        colw=[3.4, 1.5, 1.5, 1.3, 1.6, 1.2])
    bullets_slide(
        prs, "Putting it together",
        [("Three independent lines point the same way",
          "escape/rearing from the scoring (×0.08), baseline movement "
          "from the video (×0.47), and the reduction being the same size "
          "in every period"),
         ("A selective analgesic predicts none of that",
          "it should spare exploration and act only when a stimulus is "
          "applied"),
         ("So the honest description is a general reduction in activity",
          "at this dose and 10 min after injection"),
         ("What it does NOT rule out",
          "analgesia could be present as well and simply hidden. A lower "
          "dose, a later time point, and a vehicle group would separate "
          "them.")],
        "the conclusion the data supports")

    return prs


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--out", default=None)
    a = ap.parse_args()
    prs = build()
    out = a.out or os.path.join(
        OUTDIR, f"mini1p_SBI553_supplement_{date.today().isoformat()}.pptx")
    base, n = out, 1
    while True:
        try:
            prs.save(out)
            break
        except PermissionError:
            n += 1
            r, e = os.path.splitext(base)
            out = f"{r}_v{n}{e}"
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides")
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
